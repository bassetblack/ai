from fastapi import FastAPI, UploadFile, File, HTTPException, Form,Request,Body
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
import uvicorn
import json
import logging
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from PIL import Image
import pytesseract
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
from io import BytesIO
from pptx import Presentation
from docx import Document
from openpyxl import load_workbook
import PyPDF2
import fitz  # PyMuPDF
import os
import google.generativeai as genai
from main import get_response, get_gemini_response
from fastapi import Query

from pydantic import BaseModel

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI()
origins = [
    "http://localhost:5173",  # Frontend development server
    "http://127.0.0.1:5173",  # Alternative frontend URL
]

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

conversation_history = []


def process_file(contents: bytes, content_type: str) -> str:
    """
    Extracts text from the uploaded file based on its content type.
    Supported formats: PDF, DOCX, PPTX, XLSX, Images (using OCR).
    """
    try:
        if "pdf" in content_type:
            pdf_reader = PyPDF2.PdfReader(BytesIO(contents))
            text = "".join(page.extract_text() for page in pdf_reader.pages)
            return text.strip() or "No text extracted from PDF."
        elif "vnd.openxmlformats-officedocument.wordprocessingml.document" in content_type:
            doc = Document(BytesIO(contents))
            text = "\n".join(para.text for para in doc.paragraphs)
            return text.strip() or "No text extracted from DOCX."
        elif "vnd.openxmlformats-officedocument.presentationml.presentation" in content_type:
            presentation = Presentation(BytesIO(contents))
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip() or "No text extracted from PPTX."
        elif "vnd.openxmlformats-officedocument.spreadsheetml.sheet" in content_type:
            workbook = load_workbook(filename=BytesIO(contents))
            text = ""
            for sheet in workbook.sheetnames:
                worksheet = workbook[sheet]
                for row in worksheet.iter_rows(values_only=True):
                    text += " ".join(str(cell) for cell in row if cell) + "\n"
            return text.strip() or "No text extracted from XLSX."
        elif "image" in content_type:
            image = Image.open(BytesIO(contents))
            text = pytesseract.image_to_string(image)
            return text.strip() or "No text extracted from image."
        else:
            return "Unsupported file type. Please upload a PDF, DOCX, PPTX, XLSX, or image file."
    except Exception as e:
        logger.error(f"Error processing file: {e}")
        return f"Error extracting text: {e}"

@app.post("/chat")
async def chat_endpoint(chat_request: dict):
    try:
        conversation_history = chat_request.get("conversation_history", [])
        question = chat_request.get("question")

        if not isinstance(conversation_history, list):
            raise HTTPException(status_code=400, detail="Conversation history must be a list of messages.")
        if not question:
            raise HTTPException(status_code=400, detail="Question is required.")

        prompt = ""
        for message in conversation_history:
            if not all(key in message for key in ["sender_type", "text"]):
                raise HTTPException(status_code=400, detail="Each message must contain 'sender_type' and 'text'.")
            sender = "User" if message["sender_type"] == "user" else "AI"
            prompt += f"{sender}: {message['text']}\n"
        prompt += f"User: {question}\nAI:"

        response = get_response(prompt)
        return {"answer": response, "sender_type": "ai"}
    except Exception as e:
        logger.error(f"Error processing chat request: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal Server Error")
@app.post("/upload")
async def upload_file(file: UploadFile = File(...), instruction: str = Form(None), conversation_history: str = Form(None)):
    try:
        # Parse conversation history
        conversation_history = json.loads(conversation_history) if conversation_history else []
        
        # Read file contents
        contents = await file.read()
        file_name = file.filename  # Access the filename using .filename
        extracted_text = process_file(contents, file.content_type)
        
        # Build the prompt for the AI
        prompt = ""
        for message in conversation_history:
            sender = "User" if message["sender_type"] == "user" else "AI"
            prompt += f"{sender}: {message['text']}\n"
        
        # Add the instruction and extracted text to the prompt
        if instruction:
            prompt += f"User: {instruction}\nExtracted Text: {extracted_text}\nAI:"
        else:
            prompt += f"Extracted Text: {extracted_text}\nAI:"
        
        # Generate AI response
        response = get_gemini_response(prompt)
        
        # Save the file name and instruction as part of the user message
        user_message_text = f"📂 {file_name}\n"
        if instruction:
            user_message_text += f"\n{instruction}"
        
        return {
            "result": response,
            "user_message": user_message_text,  # Return the formatted user message for saving in the DB
        }
    except Exception as e:
        logger.error(f"Error processing upload request: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal Server Error")

class ExerciseRequest(BaseModel):
    prompt: str  # Topic or concept for generating the exercise

@app.post("/exercise")
async def generate_exercise(exercise_request: ExerciseRequest):
    try:
        # Extract the prompt from the request body
        prompt = exercise_request.prompt

        if not prompt.strip():
            raise HTTPException(status_code=400, detail="Prompt is required.")

        # Specialized prompt for exercise mode
        ai_prompt = (
            f"Generate a single question based on the following topic:\n"
            f"'{prompt}'\n"
            f"Do not include the correct answer in the response."
        )

        # Log the AI prompt for debugging
        logger.info("Exercise prompt sent to Gemini:\n" + ai_prompt)

        # Call the AI model to generate the exercise
        response = get_gemini_response1(ai_prompt)
        logger.info(f"AI Response: {response}")

        # Check if the response is plain text (not JSON)
        if isinstance(response, str):
            # Treat the plain text as the question
            question = response.strip()
        else:
            # Parse the AI response as JSON
            try:
                parsed_response = json.loads(response)
                question = parsed_response.get("question")
                if not question:
                    raise ValueError("AI response does not contain a valid question.")
            except (json.JSONDecodeError, ValueError) as e:
                logger.error(f"Invalid AI response: {response}")
                raise HTTPException(status_code=500, detail="Failed to parse AI response.")

        # Return the generated question
        return {"question": question}

    except HTTPException as e:
        # Re-raise HTTP exceptions (e.g., 400 Bad Request)
        raise e
    except Exception as e:
        # Log and handle all other exceptions
        logger.error(f"Error generating exercise: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal Server Error")

class ValidateAnswerRequest(BaseModel):
    user_answer: str  # User's answer to the question
    question: str     # The generated question

@app.post("/validate-answer")
async def validate_answer(validate_request: ValidateAnswerRequest):
    try:
        # Extract the user's answer and the question from the request body
        user_answer = validate_request.user_answer.strip()
        question = validate_request.question.strip()

        if not user_answer or not question:
            raise HTTPException(status_code=400, detail="Both user_answer and question are required.")

        # Specialized prompt for validating the user's answer
        ai_prompt = (
            f"Validate the user's answer for the following question:\n"
            f"Question: {question}\n"
            f"User's Answer: {user_answer}\n"
            f"Return a JSON object with the following structure:\n"
            f"{{'isCorrect': true or false, 'correctAnswer': 'The correct summerized answer'}}\n"
            f"Do not include any additional text or explanations outside the JSON object."
        )

        logger.info("Validation prompt sent to Gemini:\n" + ai_prompt)

        # Call the AI model to validate the answer
        response = get_gemini_response(ai_prompt)
        logger.info(f"AI Response: {response}")  # Log the raw AI response

        # Clean up the AI response (remove wrapping backticks and `json` prefix)
        cleaned_response = response.strip()
        if cleaned_response.startswith("```json") and cleaned_response.endswith("```"):
            cleaned_response = cleaned_response[len("```json"):].rstrip("```").strip()

        # Parse the AI response as JSON
        try:
            parsed_response = json.loads(cleaned_response)
            is_correct = parsed_response.get("isCorrect")
            correct_answer = parsed_response.get("correctAnswer")

            if is_correct is None or not isinstance(is_correct, bool):
                raise ValueError("'isCorrect' must be a boolean value.")
            if not correct_answer or not isinstance(correct_answer, str):
                raise ValueError("'correctAnswer' must be a non-empty string.")
        except (json.JSONDecodeError, ValueError) as e:
            logger.error(f"Invalid AI response: {response}")
            raise HTTPException(status_code=500, detail="Failed to parse AI response.")

        # Return the validation result and the correct answer
        return {
            "isCorrect": is_correct,
            "correctAnswer": correct_answer,
        }

    except HTTPException as e:
        # Re-raise HTTP exceptions (e.g., 400 Bad Request)
        raise e
    except Exception as e:
        # Log and handle all other exceptions
        logger.error(f"Error validating answer: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal Server Error")
# Gemini API configuration
generation_config = {
    "temperature": 0.9,
    "top_p": 0.9,
    "top_k": 40,
    "max_output_tokens": 16000,
    "response_mime_type": "text/plain",
}

model = genai.GenerativeModel(
    model_name="gemini-2.0-pro-exp-02-05",
    generation_config=generation_config,
    system_instruction=(
        "You are an AI assistant specialized in IT-related topics. Your primary role is to provide accurate and detailed answers "
        "to questions within the IT domain, including but not limited to programming, networking, cybersecurity, hardware, software, "
        "cloud computing, and troubleshooting. For non-IT topics, respond normally to general greetings, personal questions, or casual conversation. "
        "Answer math-related questions or perform calculations as needed. For all other subjects outside the IT domain, respond with: 'This is outside my dataset'. "
        "Ensure your responses are concise, professional, and tailored to the user's needs. If a question is unclear, ask for clarification before proceeding."
    ),
)

class QuizRequest(BaseModel):
    prompt: str

import re

def get_gemini_response1(prompt: str) -> str:
    try:
        response = model.generate_content(prompt)
        raw_text = response.text.strip()
        
        # Remove backticks and possible "json" label
        cleaned_text = re.sub(r"```(json)?", "", raw_text).strip()
        
        return cleaned_text
    except Exception as e:
        logger.error(f"Error in AI response: {e}")
        return ""


@app.post("/quiz")
async def generate_quiz(request: QuizRequest):
    try:
        prompt = request.prompt
        logger.info(f"Generating quiz for prompt: {prompt}")

        ai_prompt = (
            f"Generate a set of multiple-choice questions based on the following topic:\n"
            f"'{prompt}'\n"
            f" no commentary or anyhting and in answers rewrite the answer dont use alphabets for the answer or numbers rewrite them Return ONLY valid JSON in this format:\n"
            f'{{"questions": [{{"question": "example question", "options": ["A", "B", "C", "D"], "answer": "A"}}]}}'
        )

        response = get_gemini_response1(ai_prompt)

        logger.info(f"Raw AI response: {response}")  # Add this line to check response

        if not response:
            raise HTTPException(status_code=500, detail="No response from AI service")

        cleaned_response = response.strip().lstrip("'''json").rstrip("'''")

        try:
            parsed_response = json.loads(cleaned_response)
        except json.JSONDecodeError:
            raise HTTPException(status_code=500, detail="Invalid JSON response from AI service")

        questions = parsed_response.get("questions", [])
        if not questions:
            raise HTTPException(status_code=400, detail="Failed to generate valid quiz questions.")

        return {"quiz": json.dumps({"questions": questions}, ensure_ascii=False)}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"An error occurred: {str(e)}")
class ConclusionRequest(BaseModel):
    prompt: str  # Input text to generate a conclusion from

@app.post("/conclusion")
async def generate_conclusion(request: ConclusionRequest):
    try:
        prompt = request.prompt.strip()
        
        if not prompt:
            raise HTTPException(status_code=400, detail="Prompt is required.")

        # Formulate a structured AI prompt
        ai_prompt = (
            f"Based on the following content, generate a concise about the weak points and strong based on their answers to the questions, well-structured conclusion:\n"
            f"'{prompt}'\n"
            f"Ensure the conclusion is clear, impactful, and summarizes key points effectively and help them to improve. no need to explain anything in wrong or correct only show the weak points and strong based on their answers to the questions and help them to improve"
        )

        # Log the AI prompt
        logger.info(f"Generating conclusion for prompt: {prompt}")

        # Generate response from Gemini AI
        response = get_gemini_response1(ai_prompt)
        logger.info(f"AI Conclusion Response: {response}")

        if not response:
            raise HTTPException(status_code=500, detail="Failed to generate conclusion.")

        return {"conclusion": response.strip()}
    
    except HTTPException as e:
        raise e
    except Exception as e:
        logger.error(f"Error generating conclusion: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal Server Error")

if __name__ == "__main__":
    uvicorn.run("api:app", host="0.0.0.0", port=8000, reload=True)


